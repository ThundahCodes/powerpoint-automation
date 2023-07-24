import os
import time
import PySimpleGUI as sg
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

#For Loop on all processes, to check if Powerpoint is running (Made for local testing only)
"""
     for process in psutil.process_iter(['pid', 'name']):
        if program_name.lower() in process.info['name'].lower():
            try:
                pid = process.info['pid']
                os.system(f"taskkill /F /PID {pid}")  # Forcefully terminate the process
                print(f"Successfully closed '{program_name}'.")
            except psutil.NoSuchProcess:
                print(f"Failed to close '{program_name}'. Process not found.")
"""

time.sleep(0.5)
def datenSammel(lieferungDatum, ladeZeit, lieferteFirma, landFirma, rechnerAdd):

    # Presentation creation
    if os.path.isdir(fr""):
        pres = Presentation(fr"")
    else:
        pres = Presentation()

    

    # Delete all other PPTX Slides created (You need this, because when you input the Name of the PPTX as a parameter in Presentation("pptxname"), it will always create a new slide)
    for i in range(len(pres.slides)-1, -1, -1): 
        rId = pres.slides._sldIdLst[i].rId
        pres.part.drop_rel(rId)
        del pres.slides._sldIdLst[i]

    # Slide Creation, Insert a Layout in the Array
    slide_layout = pres.slide_layouts[0]
    slide = pres.slides.add_slide(slide_layout)
    pic = slide.shapes.add_picture(r"", 0, 0, width=pres.slide_width, height=pres.slide_height)
    slide.shapes._spTree.remove(pic._element)
    slide.shapes._spTree.insert(2, pic._element)
    # Title and subtitle are added with the text
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = f"{lieferteFirma} {landFirma}"
    subtitle.text = f"Abholung am {lieferungDatum} von {ladeZeit}"

    # Textbox creation with the cardinal directions set in inches
    left = Inches(6)
    height = Inches(5.5)
    top = Inches(7)
    width = Inches(5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = f"{rechnerAdd}"

    # Font and color changes in title and subtitle
    title_para = slide.shapes.title.text_frame.paragraphs[0]
    subtitle_para = slide.placeholders[1].text_frame.paragraphs[0]
    title_para.font.name = "Arial Black"
    title_para.font.size = Pt(70)
    subtitle_para.font.name = "Arial"
    subtitle_para.font.size = Pt(46)
    subtitle_para.font.color.rgb = RGBColor(255, 0, 0)
    tf.paragraphs[0].font.name = "Arial Black"
    tf.paragraphs[0].font.size = Pt(24)

    if os.path.isdir(fr""):
        pass
    else:
        os.mkdir(fr"")

    # Save the Program 
    pres.save(fr'')

#Create a GUI to insert the data into the presentation via the aforedefined Function
def main():

    sg.theme("DarkBlue3")
    layout = [
        [sg.Text("Geben Sie das Lieferungs Datum: "), sg.Input(key="lieferungDatum")],
        [sg.Text("Zeit / Ladehöhe: "), sg.Input(key="ladeZeit")],
        [sg.Text("Firma: "), sg.Input(key="lieferteFirma")],
        [sg.Text("Land: "), sg.Input(key="landFirma")],
        [sg.Text("Rechner Adresse: "), sg.Input(key="rechnerAdd")],
        [sg.Button("Run")]
    ]

    window = sg.Window("PowerPoint Generator", layout, resizable=True)

    while True:
        event, values = window.read()
        if event == sg.WIN_CLOSED:
            break
        elif event == "Run":
            datenSammel(
                values["lieferungDatum"],
                values["ladeZeit"],
                values["lieferteFirma"],
                values["landFirma"],
                values["rechnerAdd"]
            )
        sg.popup("Presentation generated successfully!", title="Success")
        window.close()
            
if __name__ == "__main__":
    main()
